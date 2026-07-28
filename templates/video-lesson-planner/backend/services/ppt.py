"""PPT generation from a lesson script using python-pptx."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

AUTHOR = "manjunath kalburgi"

# Dark theme palette
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
PRIMARY = RGBColor(0x3B, 0x82, 0xF6)
ACCENT = RGBColor(0x22, 0xD3, 0xEE)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
MUTED = RGBColor(0x94, 0xA3, 0xB8)


def _style_slide(slide, prs):
    """Apply a dark background to a slide."""
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_DARK


def _add_title_bar(slide, prs):
    """Add a colored accent bar under the title area."""
    bar = slide.shapes.add_shape(1, Inches(0.5), Inches(1.35), Inches(2.0), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()


def build_pptx(topic: str, slides: list[dict], out_path: str) -> str:
    """Render a presentation from structured slide data.

    slides: list of {"title": str, "bullets": [str, ...]}
    Slide 0 is forced to a Welcome slide; the final slide is a Thank You slide.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # --- Welcome slide (always first) ---
    s = prs.slides.add_slide(blank)
    _style_slide(s, prs)
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(2.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"Welcome to {topic}"
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = LIGHT
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = f"by {AUTHOR}"
    r2.font.size = Pt(24)
    r2.font.color.rgb = ACCENT

    # --- Content slides ---
    for slide_data in slides:
        s = prs.slides.add_slide(blank)
        _style_slide(s, prs)
        title_box = s.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12.1), Inches(0.9))
        tf = title_box.text_frame
        tf.word_wrap = True
        tp = tf.paragraphs[0]
        tr = tp.add_run()
        tr.text = slide_data.get("title", "")
        tr.font.size = Pt(32)
        tr.font.bold = True
        tr.font.color.rgb = PRIMARY
        _add_title_bar(s, prs)

        body = s.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.7), Inches(5.2))
        btf = body.text_frame
        btf.word_wrap = True
        for i, bullet in enumerate(slide_data.get("bullets", [])):
            bp = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            bp.space_after = Pt(12)
            br = bp.add_run()
            br.text = f"•  {bullet}"
            br.font.size = Pt(20)
            br.font.color.rgb = LIGHT

    # --- Thank you slide (always last) ---
    s = prs.slides.add_slide(blank)
    _style_slide(s, prs)
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.7), Inches(2.0))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Thank You"
    run.font.size = Pt(48)
    run.font.bold = True
    run.font.color.rgb = LIGHT
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = f"by {AUTHOR}"
    r2.font.size = Pt(26)
    r2.font.color.rgb = ACCENT

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return out_path
